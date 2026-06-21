from __future__ import annotations

import io
import re
from collections import Counter
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


EMU_PER_INCH = 914400
POINTS_PER_EMU = 72 / EMU_PER_INCH


def parse_per_slide_intents(raw_text: str | None) -> dict[int, str]:
    """
    スライドごとの意図メッセージ文字列を解析してスライド番号と意図メッセージの辞書を返す

    Args
    -----------------
    - raw_text: str | None,     「番号: メッセージ」形式の改行区切りテキスト

    Returns
    -----------------
    - intents: dict[int, str],  スライド番号をキー、意図メッセージを値とする辞書

    """
    # テキストが空の場合は空辞書を返却
    if not raw_text:
        return {}
    intents: dict[int, str] = {}
    # 行ごとに解析
    for line in raw_text.splitlines():
        line = line.strip()
        # 空行はスキップ
        if not line:
            continue
        # 「番号: メッセージ」形式にマッチするか確認
        match = re.match(r"^(\d+)\s*[:：-]\s*(.+)$", line)
        if match:
            # スライド番号と意図メッセージを辞書に追加
            intents[int(match.group(1))] = match.group(2).strip()
    return intents


def _safe_text(value: str | None) -> str:
    """
    Noneまたは空白のみの文字列を空文字列に変換する

    Args
    -----------------
    - value: str | None,    変換対象の文字列

    Returns
    -----------------
    - result: str,          空白を除去した文字列（Noneの場合は空文字列）

    """
    return (value or "").strip()


def _shape_type_name(shape: Any) -> str:
    """
    シェイプの種別名を文字列で返す

    Args
    -----------------
    - shape: Any,       PPTXのシェイプオブジェクト

    Returns
    -----------------
    - name: str,        シェイプ種別名（取得できない場合は文字列変換した値）

    """
    try:
        return MSO_SHAPE_TYPE(shape.shape_type).name
    except Exception:
        # 種別が取得できない場合は文字列に変換して返却
        return str(getattr(shape, "shape_type", "UNKNOWN"))


def _extract_text_frame(shape: Any) -> dict[str, Any] | None:
    """
    シェイプからテキストフレームの情報を抽出して返す

    Args
    -----------------
    - shape: Any,                   PPTXのシェイプオブジェクト

    Returns
    -----------------
    - text_info: dict | None,       テキスト情報の辞書（テキストなしの場合はNone）

    """
    # テキストフレームを持たないシェイプはスキップ
    if not getattr(shape, "has_text_frame", False):
        return None

    # 各種集計用の入れ物を用意
    paragraphs = []
    font_sizes_pt: list[float] = []
    total_runs = 0
    raw_text_parts: list[str] = []
    bullet_count = 0

    ## 段落ごとにテキスト情報を収集
    for paragraph in shape.text_frame.paragraphs:
        # 段落のテキストを取得
        paragraph_text = "".join(run.text for run in paragraph.runs).strip() or paragraph.text.strip()
        if paragraph_text:
            raw_text_parts.append(paragraph_text)
            paragraphs.append({
                "text": paragraph_text,
                "level": getattr(paragraph, "level", 0),
                "alignment": str(getattr(paragraph, "alignment", None)),
            })
            # 2段落目以降またはインデントありの段落はバレットとしてカウント
            if len(paragraphs) > 1 or getattr(paragraph, "level", 0) > 0:
                bullet_count += 1

        # ランごとにフォントサイズを収集
        for run in paragraph.runs:
            total_runs += 1
            size = getattr(getattr(run, "font", None), "size", None)
            if size is not None:
                try:
                    font_sizes_pt.append(round(size.pt, 2))
                except Exception:
                    pass

    # テキストを結合
    text = "\n".join(raw_text_parts).strip()
    # テキストが空の場合はNoneを返却
    if not text:
        return None

    return {
        "text": text,
        "paragraphs": paragraphs,
        "bullet_count": bullet_count,
        "font_sizes_pt": font_sizes_pt,
        "explicit_font_run_count": total_runs,
    }


def _bbox(shape: Any) -> dict[str, float]:
    """
    シェイプの絶対座標（EMU単位）を辞書で返す

    Args
    -----------------
    - shape: Any,               PPTXのシェイプオブジェクト

    Returns
    -----------------
    - bbox: dict[str, float],   left, top, width, height（EMU単位）

    """
    return {
        "left": round(float(shape.left), 2),
        "top": round(float(shape.top), 2),
        "width": round(float(shape.width), 2),
        "height": round(float(shape.height), 2),
    }


def _bbox_ratios(shape: Any, slide_width: int, slide_height: int) -> dict[str, float]:
    """
    シェイプの座標をスライドサイズに対する比率で返す

    Args
    -----------------
    - shape: Any,                       PPTXのシェイプオブジェクト
    - slide_width: int,                 スライドの幅（EMU単位）
    - slide_height: int,                スライドの高さ（EMU単位）

    Returns
    -----------------
    - ratios: dict[str, float],         x_ratio, y_ratio, w_ratio, h_ratio（0.0〜1.0）

    """
    return {
        "x_ratio": round(float(shape.left) / slide_width, 4) if slide_width else 0.0,
        "y_ratio": round(float(shape.top) / slide_height, 4) if slide_height else 0.0,
        "w_ratio": round(float(shape.width) / slide_width, 4) if slide_width else 0.0,
        "h_ratio": round(float(shape.height) / slide_height, 4) if slide_height else 0.0,
    }


def _slide_heuristics(slide_data: dict[str, Any], slide_width: int, slide_height: int) -> dict[str, Any]:
    """
    スライドデータからヒューリスティックな品質フラグを生成して返す

    Args
    -----------------
    - slide_data: dict[str, Any],       解析済みのスライドデータ
    - slide_width: int,                 スライドの幅（EMU単位）
    - slide_height: int,                スライドの高さ（EMU単位）

    Returns
    -----------------
    - heuristics: dict[str, Any],       テキスト面積比率・構造・ビジュアル・コンテンツのフラグ

    """
    ## テキスト面積比率を計算
    slide_area = slide_width * slide_height if slide_width and slide_height else 1
    text_area = 0
    for block in slide_data["text_blocks"]:
        bbox = block["bbox"]
        text_area += bbox["width"] * bbox["height"]
    # テキストが占める面積の割合
    text_area_ratio = round(text_area / slide_area, 4)

    # 各種メトリクスを取得
    total_chars = slide_data["text_char_count"]
    title = slide_data.get("title") or ""
    bullet_count = slide_data["bullet_count"]
    text_box_count = len(slide_data["text_blocks"])
    avg_font_size = slide_data["font_stats"].get("average_pt")

    ## 構造上の問題フラグを生成
    structure_flags = []
    # タイトルが存在しない場合
    if not title:
        structure_flags.append("title_missing")
    # タイトルが長すぎる場合
    elif len(title) > 60:
        structure_flags.append("title_may_not_be_concise")
    # バレットが多すぎる場合
    if bullet_count >= 7:
        structure_flags.append("too_many_bullets")
    # テキストボックスが多すぎる場合
    if text_box_count >= 6:
        structure_flags.append("too_many_text_regions")

    ## ビジュアル上の問題フラグを生成
    visual_flags = []
    # テキストの占有面積が広すぎる場合
    if text_area_ratio > 0.45:
        visual_flags.append("text_dense_layout")
    # フォントサイズが小さすぎる場合
    if avg_font_size is not None and avg_font_size < 18:
        visual_flags.append("small_explicit_font_size")
    # 画像なしでテキストが多い場合
    if slide_data["shape_counts"].get("PICTURE", 0) == 0 and total_chars > 220:
        visual_flags.append("text_heavy_without_visual_anchor")

    ## コンテンツ上の問題フラグを生成
    content_flags = []
    # テキストが極端に少ない場合
    if total_chars < 20:
        content_flags.append("very_little_text")
    # タイトルなしでテキストがある場合
    if not title and total_chars > 0:
        content_flags.append("message_without_clear_heading")

    return {
        "text_area_ratio": text_area_ratio,
        "structure_flags": structure_flags,
        "visual_flags": visual_flags,
        "content_flags": content_flags,
        "candidate_review_risks": structure_flags + visual_flags + content_flags,
    }


def parse_pptx(file_bytes: bytes, filename: str, overall_intended_message: str = "", per_slide_intended_messages: str | None = None) -> dict[str, Any]:
    """
    PPTXファイルのバイト列を解析してスライド情報の辞書を返す

    Args
    -----------------
    - file_bytes: bytes,                        PPTXファイルのバイト列
    - filename: str,                            ファイル名
    - overall_intended_message: str,            プレゼンテーション全体の意図メッセージ
    - per_slide_intended_messages: str | None,  スライドごとの意図メッセージ（改行区切り）

    Returns
    -----------------
    - result: dict[str, Any],                   ファイル名・スライドサイズ・スライド一覧などの解析結果

    """
    # PPTXファイルを読み込む
    prs = Presentation(io.BytesIO(file_bytes))
    # スライドサイズを取得
    slide_width = int(prs.slide_width)
    slide_height = int(prs.slide_height)
    # スライドごとの意図メッセージを解析
    per_slide_intent_map = parse_per_slide_intents(per_slide_intended_messages)

    # 入れ物を用意
    slides: list[dict[str, Any]] = []
    deck_shape_counter: Counter[str] = Counter()

    ## スライドごとに解析
    for index, slide in enumerate(prs.slides, start=1):
        # タイトルシェイプを取得
        title_shape = slide.shapes.title if hasattr(slide.shapes, "title") else None
        title = _safe_text(getattr(title_shape, "text", "")) if title_shape is not None else ""

        # スライドデータの初期値を設定
        slide_data: dict[str, Any] = {
            "slide_number": index,
            "title": title,
            # スライド固有の意図メッセージがあれば優先し、なければ全体の意図メッセージを使用
            "intended_message": per_slide_intent_map.get(index) or overall_intended_message.strip(),
            "text_blocks": [],
            "table_count": 0,
            "chart_count": 0,
            "image_count": 0,
            "placeholder_count": 0,
            "shape_counts": {},
            "text_char_count": 0,
            "bullet_count": 0,
            "font_stats": {},
            "all_text": "",
        }

        # 集計用の入れ物を用意
        font_sizes_pt: list[float] = []
        texts_for_slide: list[str] = []
        slide_shape_counter: Counter[str] = Counter()

        ## シェイプごとに解析
        for shape in slide.shapes:
            # シェイプの種別名を取得してカウント
            shape_name = _shape_type_name(shape)
            slide_shape_counter[shape_name] += 1
            deck_shape_counter[shape_name] += 1

            # プレースホルダーの場合はカウント
            if getattr(shape, "is_placeholder", False):
                slide_data["placeholder_count"] += 1
            # テーブルの場合はカウント
            if getattr(shape, "has_table", False):
                slide_data["table_count"] += 1
            # チャートの場合はカウント
            if getattr(shape, "has_chart", False):
                slide_data["chart_count"] += 1
            # 画像の場合はカウント
            if shape_name == "PICTURE":
                slide_data["image_count"] += 1

            # テキストフレームの情報を抽出
            text_info = _extract_text_frame(shape)
            if text_info:
                texts_for_slide.append(text_info["text"])
                slide_data["bullet_count"] += text_info["bullet_count"]
                font_sizes_pt.extend(text_info["font_sizes_pt"])
                # テキストブロック情報を追加
                slide_data["text_blocks"].append({
                    "shape_name": getattr(shape, "name", ""),
                    "shape_type": shape_name,
                    "bbox": _bbox(shape),
                    "bbox_ratios": _bbox_ratios(shape, slide_width, slide_height),
                    "text": text_info["text"],
                    "paragraphs": text_info["paragraphs"],
                    "bullet_count": text_info["bullet_count"],
                })

        # スライド全体のテキストを結合して文字数を計算
        all_text = "\n\n".join(texts_for_slide).strip()
        slide_data["all_text"] = all_text
        slide_data["text_char_count"] = len(all_text.replace("\n", ""))
        slide_data["shape_counts"] = dict(slide_shape_counter)

        ## フォントサイズの統計情報を集計
        if font_sizes_pt:
            slide_data["font_stats"] = {
                "min_pt": round(min(font_sizes_pt), 2),
                "max_pt": round(max(font_sizes_pt), 2),
                "average_pt": round(sum(font_sizes_pt) / len(font_sizes_pt), 2),
                "sample_count": len(font_sizes_pt),
            }
        else:
            # フォントサイズ情報が取得できなかった場合はNoneを設定
            slide_data["font_stats"] = {
                "min_pt": None,
                "max_pt": None,
                "average_pt": None,
                "sample_count": 0,
            }

        # ヒューリスティックな品質フラグを生成して追加
        slide_data["heuristics"] = _slide_heuristics(slide_data, slide_width, slide_height)
        slides.append(slide_data)

    # レスポンス返却
    return {
        "file_name": filename,
        "slide_count": len(slides),
        "slide_size": {
            "width_emu": slide_width,
            "height_emu": slide_height,
            "width_inches": round(slide_width / EMU_PER_INCH, 2),
            "height_inches": round(slide_height / EMU_PER_INCH, 2),
        },
        "overall_intended_message": overall_intended_message.strip(),
        "per_slide_intended_messages": per_slide_intent_map,
        "deck_shape_counts": dict(deck_shape_counter),
        "slides": slides,
    }
