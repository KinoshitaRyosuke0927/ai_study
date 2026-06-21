from __future__ import annotations

import base64
import io

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation

THUMB_WIDTH = 800

# フォント候補パスリスト（OS 環境ごとに優先順位を付けて探索）
_FONT_CANDIDATES = [
    "C:\\Windows\\Fonts\\YuGothM.ttc",
    "C:\\Windows\\Fonts\\meiryo.ttc",
    "C:\\Windows\\Fonts\\msgothic.ttc",
    "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
]


def _load_font(size: int) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    """
    指定サイズのフォントを候補パスから順に探してロードする

    Args
    -----------------
    - size: int,                                            フォントサイズ（pt）

    Returns
    -----------------
    - font: ImageFont.FreeTypeFont | ImageFont.ImageFont,   ロードしたフォントオブジェクト

    """
    for path in _FONT_CANDIDATES:
        try:
            return ImageFont.truetype(path, size)
        except Exception:
            pass
    # 候補がすべて失敗した場合はデフォルトフォントを使用
    return ImageFont.load_default()


def render_slide_thumbnails(file_bytes: bytes) -> list[str]:
    """
    PPTXの各スライドをPNG画像（Base64）のリストに変換する

    Args
    -----------------
    - file_bytes: bytes,    PPTXファイルのバイト列

    Returns
    -----------------
    - results: list[str],   スライドごとの PNG 画像を Base64 エンコードした文字列リスト

    """
    # PPTXを読み込み、サムネイルサイズを計算
    prs = Presentation(io.BytesIO(file_bytes))
    sw = int(prs.slide_width) or 1
    sh = int(prs.slide_height) or 1
    thumb_h = int(THUMB_WIDTH * sh / sw)
    scale = THUMB_WIDTH / sw

    title_font = _load_font(16)
    body_font = _load_font(12)

    results: list[str] = []
    ## スライドごとにサムネイルを生成
    for slide in prs.slides:
        img = Image.new("RGB", (THUMB_WIDTH, thumb_h), "#f8f9fa")
        draw = ImageDraw.Draw(img)

        title_shape = slide.shapes.title

        # シェイプごとにテキストを描画
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue

            x0 = int(shape.left * scale)
            y0 = int(shape.top * scale)
            x1 = x0 + int(shape.width * scale)
            y1 = y0 + int(shape.height * scale)
            x0, x1 = max(0, x0), min(THUMB_WIDTH, x1)
            y0, y1 = max(0, y0), min(thumb_h, y1)

            is_title = shape is title_shape
            bg_color = "#ffffff" if is_title else "#ffffff"
            draw.rectangle([x0, y0, x1, y1], fill=bg_color, outline="#d0d7de", width=1)

            font = title_font if is_title else body_font
            line_h = 18 if is_title else 14
            cursor_y = y0 + 4

            # 段落ごとにテキストをクリップして描画
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                if cursor_y + line_h > y1 - 2:
                    break
                clip_x = x1 - x0 - 8
                truncated = _clip_text(text, font, clip_x)
                color = "#0052cc" if is_title else "#24292f"
                draw.text((x0 + 4, cursor_y), truncated, font=font, fill=color)
                cursor_y += line_h + 2

        # PNG に変換して Base64 エンコード
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        results.append(base64.b64encode(buf.getvalue()).decode())

    return results


def _clip_text(text: str, font: ImageFont.FreeTypeFont | ImageFont.ImageFont, max_px: int) -> str:
    """
    テキストを指定ピクセル幅に収まるようにクリップして返す

    Args
    -----------------
    - text: str,                                            クリップ対象のテキスト
    - font: ImageFont.FreeTypeFont | ImageFont.ImageFont,   テキスト描画に使用するフォント
    - max_px: int,                                          最大幅（ピクセル）

    Returns
    -----------------
    - clipped: str,     クリップ後のテキスト（省略の場合は末尾に「…」を付与）

    """
    try:
        bbox = font.getbbox(text)
        if bbox[2] <= max_px:
            return text
        # 末尾から1文字ずつ削りながら収まる長さを探す
        for i in range(len(text), 0, -1):
            t = text[:i] + "…"
            if font.getbbox(t)[2] <= max_px:
                return t
        return "…"
    except Exception:
        return text
